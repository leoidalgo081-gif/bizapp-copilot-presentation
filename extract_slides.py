import os
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_pptx_data(file_path):
    print(f"Extracting: {file_path}")
    prs = Presentation(file_path)
    slides_data = []
    
    # Create directory for images extracted from this presentation
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    img_dir = os.path.join("extracted_images", base_name)
    os.makedirs(img_dir, exist_ok=True)
    
    img_counter = 0
    for i, slide in enumerate(prs.slides):
        slide_title = ""
        # Check if slide has a title shape
        try:
            if slide.shapes.title:
                slide_title = slide.shapes.title.text
        except Exception:
            pass
            
        slide_text_elements = []
        images_in_slide = []
        
        for shape in slide.shapes:
            # Extract text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text_elements.append(text)
            
            # Extract tables
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_text_elements.append(" | ".join(row_text))
            
            # Extract images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    img_name = f"slide_{i+1}_img_{img_counter}.{ext}"
                    img_path = os.path.join(img_dir, img_name)
                    # We store it using forward slashes for cross-platform ease
                    web_img_path = f"extracted_images/{base_name}/{img_name}"
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    images_in_slide.append(web_img_path)
                    img_counter += 1
                except Exception as e:
                    print(f"Error extracting image from slide {i+1}: {e}")
                
        # Consolidate slide text
        slides_data.append({
            "slide_number": i + 1,
            "title": slide_title,
            "text": "\n".join(slide_text_elements),
            "images": images_in_slide
        })
        
    return slides_data

def main():
    files = ["NOVO.pptx", "Template Apresentação com Copilot - Bizapp 2026.pptx", "Treinamento 2.pptx"]
    all_data = {}
    for file in files:
        if os.path.exists(file):
            try:
                all_data[file] = extract_pptx_data(file)
            except Exception as e:
                print(f"Error extracting {file}: {e}")
        else:
            print(f"File not found: {file}")
            
    with open("slides_summary.json", "w", encoding="utf-8") as f:
        json.dump(all_data, f, indent=4, ensure_ascii=False)
    print("Done! Extracted data written to slides_summary.json")

if __name__ == "__main__":
    main()
